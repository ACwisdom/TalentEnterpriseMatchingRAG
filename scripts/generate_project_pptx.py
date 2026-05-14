"""
生成根目录项目介绍 PPT（16:9）。

依赖: pip install python-pptx

用法（仓库根目录）:
  python scripts/generate_project_pptx.py
"""

from __future__ import annotations

from pathlib import Path

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import MSO_ANCHOR, PP_ALIGN
from pptx.util import Inches, Pt


# 配色：深蓝主色 + 青绿强调（人才/政企感）
C_BG_DARK = RGBColor(0x1A, 0x1F, 0x3C)
C_BG_LIGHT = RGBColor(0xF4, 0xF6, 0xFA)
C_ACCENT = RGBColor(0x02, 0x80, 0x90)
C_TEXT_ON_DARK = RGBColor(0xFF, 0xFF, 0xFF)
C_TEXT_MUTED = RGBColor(0xC8, 0xD0, 0xE0)
C_TEXT_BODY = RGBColor(0x22, 0x28, 0x38)


def _blank_slide(prs: Presentation):
    for i in (6, 5, 1):
        if i < len(prs.slide_layouts):
            return prs.slides.add_slide(prs.slide_layouts[i])
    return prs.slides.add_slide(prs.slide_layouts[0])


def _slide_bg_rect(slide, prs: Presentation, rgb: RGBColor) -> None:
    shp = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE,
        0,
        0,
        prs.slide_width,
        prs.slide_height,
    )
    shp.fill.solid()
    shp.fill.fore_color.rgb = rgb
    shp.line.fill.background()


def _add_textbox(
    slide,
    left,
    top,
    width,
    height,
    text: str,
    *,
    size_pt: int = 18,
    bold: bool = False,
    color: RGBColor = C_TEXT_BODY,
    align=PP_ALIGN.LEFT,
    valign=MSO_ANCHOR.TOP,
) -> None:
    box = slide.shapes.add_textbox(left, top, width, height)
    tf = box.text_frame
    tf.word_wrap = True
    tf.vertical_anchor = valign
    tf.clear()
    p = tf.paragraphs[0]
    p.text = text
    p.alignment = align
    p.font.size = Pt(size_pt)
    p.font.bold = bold
    p.font.color.rgb = color
    p.font.name = "Microsoft YaHei"


def _accent_bar(slide, prs: Presentation, top: float) -> None:
    bar = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE,
        Inches(0.55),
        Inches(top),
        Inches(0.12),
        Inches(0.55),
    )
    bar.fill.solid()
    bar.fill.fore_color.rgb = C_ACCENT
    bar.line.fill.background()


def main() -> None:
    root = Path(__file__).resolve().parent.parent
    out = root / "TalentEnterpriseMatchingRAG_project_deck.pptx"

    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)

    # --- 1 封面 ---
    s = _blank_slide(prs)
    _slide_bg_rect(s, prs, C_BG_DARK)
    _add_textbox(
        s,
        Inches(0.85),
        Inches(2.1),
        Inches(11.5),
        Inches(1.2),
        "TalentEnterpriseMatchingRAG",
        size_pt=40,
        bold=True,
        color=C_TEXT_ON_DARK,
    )
    _add_textbox(
        s,
        Inches(0.85),
        Inches(3.15),
        Inches(11.5),
        Inches(0.9),
        "人才—企业智能匹配与报告生成（RAG）",
        size_pt=22,
        color=C_TEXT_MUTED,
    )
    _add_textbox(
        s,
        Inches(0.85),
        Inches(6.35),
        Inches(11.5),
        Inches(0.45),
        "企业库检索 · MiMo 报告 · Word 导出 · Web 流式体验",
        size_pt=12,
        color=C_TEXT_MUTED,
    )

    # --- 2 一句话 ---
    s = _blank_slide(prs)
    _slide_bg_rect(s, prs, C_BG_LIGHT)
    _accent_bar(s, prs, 0.65)
    _add_textbox(s, Inches(0.9), Inches(0.55), Inches(11.5), Inches(0.7), "项目定位", size_pt=32, bold=True, color=C_BG_DARK)
    _add_textbox(
        s,
        Inches(0.9),
        Inches(1.45),
        Inches(11.5),
        Inches(2.2),
        "以结构化企业库与向量索引为基础，面向高层次人才简历，完成「召回—重排—解释—成稿」闭环；"
        "可选大模型生成匹配报告与画像，并导出 Word，支撑咨询与政企材料起草。",
        size_pt=20,
        color=C_TEXT_BODY,
    )
    # 右侧大号数字强调
    stat = s.shapes.add_shape(MSO_SHAPE.OVAL, Inches(10.2), Inches(2.0), Inches(2.4), Inches(2.4))
    stat.fill.solid()
    stat.fill.fore_color.rgb = C_ACCENT
    stat.line.fill.background()
    _add_textbox(s, Inches(10.35), Inches(2.75), Inches(2.1), Inches(0.9), "RAG", size_pt=36, bold=True, color=C_TEXT_ON_DARK, align=PP_ALIGN.CENTER)

    # --- 3 典型场景 ---
    s = _blank_slide(prs)
    _slide_bg_rect(s, prs, C_BG_LIGHT)
    _accent_bar(s, prs, 0.65)
    _add_textbox(s, Inches(0.9), Inches(0.55), Inches(11.5), Inches(0.7), "典型场景", size_pt=32, bold=True, color=C_BG_DARK)
    bullets = (
        "引才咨询：博士/高层次人才简历与企业岗位或园区政策快速对齐\n"
        "政企对接：批量初筛、形成可读报告，减少手工检索与复制粘贴\n"
        "内部评审：统一检索口径与重排规则，便于复核与留痕（结合业务系统扩展）"
    )
    _add_textbox(s, Inches(0.9), Inches(1.4), Inches(7.2), Inches(3.5), bullets, size_pt=17, color=C_TEXT_BODY)
    # 右侧卡片
    card = s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(8.35), Inches(1.35), Inches(4.35), Inches(4.8))
    card.fill.solid()
    card.fill.fore_color.rgb = RGBColor(0xFF, 0xFF, 0xFF)
    card.line.color.rgb = RGBColor(0xD8, 0xDE, 0xE8)
    _add_textbox(
        s,
        Inches(8.55),
        Inches(1.55),
        Inches(3.95),
        Inches(4.3),
        "输入\n• 简历 PDF / DOCX / TXT\n• 企业 Excel 库\n\n输出\n• 匹配列表 + 解释\n• 报告 / Word\n• Web 流式展示",
        size_pt=15,
        color=C_TEXT_BODY,
    )

    # --- 4 能力地图 ---
    s = _blank_slide(prs)
    _slide_bg_rect(s, prs, C_BG_DARK)
    _add_textbox(s, Inches(0.85), Inches(0.55), Inches(11.5), Inches(0.75), "核心能力一览", size_pt=32, bold=True, color=C_TEXT_ON_DARK)
    rows = [
        ("企业检索", "关键词 / 向量 / 混合（RRF）"),
        ("向量索引", "Chroma + sentence-transformers"),
        ("简历解析", "PDF、DOCX、TXT"),
        ("匹配重排", "宽召回 + CrossEncoder + 规则 + 可选 LTR"),
        ("大模型", "小米 MiMo 报告与画像"),
        ("Web", "FastAPI · SSE 流式 · /demo"),
        ("扩展", "Java REST v1 · LangGraph 工具层"),
    ]
    y = 1.35
    for i, (t1, t2) in enumerate(rows):
        oy = y + i * 0.78
        dot = s.shapes.add_shape(MSO_SHAPE.OVAL, Inches(0.85), Inches(oy), Inches(0.22), Inches(0.22))
        dot.fill.solid()
        dot.fill.fore_color.rgb = C_ACCENT
        dot.line.fill.background()
        _add_textbox(s, Inches(1.2), Inches(oy - 0.02), Inches(3.1), Inches(0.45), t1, size_pt=16, bold=True, color=C_TEXT_ON_DARK)
        _add_textbox(s, Inches(4.45), Inches(oy - 0.02), Inches(8.0), Inches(0.45), t2, size_pt=15, color=C_TEXT_MUTED)

    # --- 5 技术架构（示意）---
    s = _blank_slide(prs)
    _slide_bg_rect(s, prs, C_BG_LIGHT)
    _accent_bar(s, prs, 0.65)
    _add_textbox(s, Inches(0.9), Inches(0.55), Inches(11.5), Inches(0.7), "技术架构（示意）", size_pt=32, bold=True, color=C_BG_DARK)
    layers = [
        ("表现层", "web/main.py · 静态页 · SSE"),
        ("编排与业务", "src/pipeline.py · 约束与重排"),
        ("检索层", "src/enterprise_search.py · 向量 / 混合"),
        ("数据与索引", "data/ · Chroma 持久化"),
        ("模型服务", "MiMo API（可选）"),
    ]
    bx = Inches(0.95)
    bw = Inches(11.4)
    bh = Inches(0.78)
    y0 = 1.25
    for i, (title, sub) in enumerate(layers):
        top_in = y0 + i * 0.88
        r = s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, bx, Inches(top_in), bw, bh)
        r.fill.solid()
        r.fill.fore_color.rgb = RGBColor(0xFF, 0xFF, 0xFF)
        r.line.color.rgb = RGBColor(0xC8, 0xD0, 0xE0)
        _add_textbox(s, Inches(1.15), Inches(top_in + 0.12), Inches(3.0), Inches(0.5), title, size_pt=17, bold=True, color=C_ACCENT)
        _add_textbox(s, Inches(4.2), Inches(top_in + 0.12), Inches(7.8), Inches(0.5), sub, size_pt=15, color=C_TEXT_BODY)

    # --- 6 端到端流程 ---
    s = _blank_slide(prs)
    _slide_bg_rect(s, prs, C_BG_LIGHT)
    _accent_bar(s, prs, 0.65)
    _add_textbox(s, Inches(0.9), Inches(0.55), Inches(11.5), Inches(0.7), "端到端流程", size_pt=32, bold=True, color=C_BG_DARK)
    steps = ["解析简历", "构建查询", "宽召回", "重排与解释", "组装提示词", "生成报告", "导出 Word"]
    x0 = 0.55
    for i, st in enumerate(steps):
        xi = x0 + i * 1.72
        circ = s.shapes.add_shape(MSO_SHAPE.OVAL, Inches(xi), Inches(2.0), Inches(0.55), Inches(0.55))
        circ.fill.solid()
        circ.fill.fore_color.rgb = C_ACCENT if i < len(steps) - 1 else C_BG_DARK
        circ.line.fill.background()
        _add_textbox(
            s,
            Inches(xi - 0.05),
            Inches(2.08),
            Inches(0.65),
            Inches(0.45),
            str(i + 1),
            size_pt=14,
            bold=True,
            color=C_TEXT_ON_DARK,
            align=PP_ALIGN.CENTER,
        )
        if i < len(steps) - 1:
            arr = s.shapes.add_shape(MSO_SHAPE.RIGHT_ARROW, Inches(xi + 0.62), Inches(2.18), Inches(0.45), Inches(0.22))
            arr.fill.solid()
            arr.fill.fore_color.rgb = RGBColor(0xAA, 0xB4, 0xC5)
            arr.line.fill.background()
        _add_textbox(s, Inches(xi - 0.35), Inches(2.85), Inches(1.55), Inches(0.85), st, size_pt=12, bold=True, color=C_TEXT_BODY, align=PP_ALIGN.CENTER)
    _add_textbox(
        s,
        Inches(0.9),
        Inches(4.35),
        Inches(11.5),
        Inches(1.5),
        "说明：检索模式、重排开关与 LTR 等由环境变量与 src/config.py 控制；详见 docs/match-rerank.md。",
        size_pt=14,
        color=C_TEXT_BODY,
    )

    # --- 7 部署与运维 ---
    s = _blank_slide(prs)
    _slide_bg_rect(s, prs, C_BG_DARK)
    _add_textbox(s, Inches(0.85), Inches(0.55), Inches(11.5), Inches(0.75), "部署与运维要点", size_pt=32, bold=True, color=C_TEXT_ON_DARK)
    txt = (
        "• Docker / Fly.io：见 docs/deploy-public.md\n"
        "• 公网 SSE：反向代理需关闭缓冲、合理超时\n"
        "• 密钥：MIMO_API_KEY、APP_API_KEY（Web）\n"
        "• 可选 Java：`java/recruitment-api` + `langgraph_tools`（docs/java-langgraph-integration.md）"
    )
    _add_textbox(s, Inches(0.9), Inches(1.45), Inches(11.2), Inches(4.5), txt, size_pt=18, color=C_TEXT_MUTED)

    # --- 8 仓库与文档 ---
    s = _blank_slide(prs)
    _slide_bg_rect(s, prs, C_BG_LIGHT)
    _accent_bar(s, prs, 0.65)
    _add_textbox(s, Inches(0.9), Inches(0.55), Inches(11.5), Inches(0.7), "仓库与文档索引", size_pt=32, bold=True, color=C_BG_DARK)
    doclist = (
        "README.md — 总览与命令\n"
        "docs/项目介绍.md — 约 300 字简介\n"
        "docs/deploy-public.md — 公网部署\n"
        "docs/match-rerank.md — 匹配与重排\n"
        "docs/java-langgraph-integration.md — Java 与 LangGraph\n"
        "TalentEnterpriseMatchingRAG_project_brief_feishu.docx — 飞书用 Word 说明（根目录）"
    )
    _add_textbox(s, Inches(0.9), Inches(1.35), Inches(11.2), Inches(4.8), doclist, size_pt=16, color=C_TEXT_BODY)

    # --- 9 结语 ---
    s = _blank_slide(prs)
    _slide_bg_rect(s, prs, C_BG_DARK)
    _add_textbox(
        s,
        Inches(0.9),
        Inches(2.6),
        Inches(11.5),
        Inches(1.0),
        "谢谢",
        size_pt=44,
        bold=True,
        color=C_TEXT_ON_DARK,
        align=PP_ALIGN.CENTER,
    )
    _add_textbox(
        s,
        Inches(0.9),
        Inches(3.75),
        Inches(11.5),
        Inches(0.8),
        "欢迎提问与交流",
        size_pt=22,
        color=C_TEXT_MUTED,
        align=PP_ALIGN.CENTER,
    )

    prs.save(out)
    print(f"Wrote: {out}")


if __name__ == "__main__":
    main()
